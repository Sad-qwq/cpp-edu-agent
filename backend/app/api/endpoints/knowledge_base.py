from datetime import datetime
from pathlib import Path
from typing import Any, Optional

from fastapi import APIRouter, Depends, HTTPException
from sqlalchemy.ext.asyncio import AsyncSession
from sqlmodel import delete, func, select

from app.ai.retrieval.chunking import estimate_token_count, split_text
from app.api import deps
from app.db.session import get_session
from app.models.ai_question_generation import (
    KnowledgeChunk,
    KnowledgeDocument,
    KnowledgeParseStatus,
    KnowledgeSourceType,
)
from app.models.classroom import Classroom
from app.models.material import Material
from app.models.user import User, UserRole
from app.schemas.ai_question_generation import (
    KnowledgeDocumentRead,
    KnowledgeSearchResponse,
)

router = APIRouter()

BACKEND_ROOT = Path(__file__).resolve().parents[3]


async def _ensure_material_ingest_permission(session: AsyncSession, material: Material, current_user: User) -> Classroom:
    classroom = await session.get(Classroom, material.class_id)
    if not classroom:
        raise HTTPException(status_code=404, detail="Classroom not found")

    if current_user.role == UserRole.ADMIN:
        return classroom

    if current_user.role == UserRole.TEACHER and classroom.teacher_id == current_user.id:
        return classroom

    raise HTTPException(status_code=403, detail="Not enough permissions")


def _resolve_material_file_path(file_url: str) -> Path:
    relative_path = file_url.replace("/static/", "uploads/").lstrip("/")
    return BACKEND_ROOT / relative_path


def _extract_text_from_material(material: Material, file_path: Path) -> tuple[str, dict[str, Any]]:
    suffix = file_path.suffix.lower()
    metadata = {
        "material_id": material.id,
        "material_type": material.file_type,
        "extract_mode": "fallback",
    }

    if suffix in {".txt", ".md", ".py", ".cpp", ".c", ".h", ".hpp", ".json"} and file_path.exists():
        text = file_path.read_text(encoding="utf-8", errors="ignore")
        metadata["extract_mode"] = "direct_text"
        return text, metadata

    if suffix == ".pdf" and file_path.exists():
        from pypdf import PdfReader

        reader = PdfReader(str(file_path))
        text = "\n".join((page.extract_text() or "") for page in reader.pages)
        metadata["extract_mode"] = "pdf"
        return text, metadata

    if suffix in {".ppt", ".pptx"} and file_path.exists():
        from pptx import Presentation

        presentation = Presentation(str(file_path))
        texts: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        metadata["extract_mode"] = "ppt"
        return "\n".join(texts), metadata

    fallback_parts = [material.title]
    if material.description:
        fallback_parts.append(material.description)
    fallback_parts.append(f"资料类型: {material.file_type}")
    if material.file_url:
        fallback_parts.append(f"文件地址: {material.file_url}")
    return "\n".join(fallback_parts), metadata


@router.post("/materials/{material_id}/ingest", response_model=KnowledgeDocumentRead)
async def ingest_material_to_knowledge_base(
    material_id: int,
    *,
    session: AsyncSession = Depends(get_session),
    current_user: User = Depends(deps.get_current_active_user),
) -> Any:
    material = await session.get(Material, material_id)
    if not material:
        raise HTTPException(status_code=404, detail="Material not found")

    classroom = await _ensure_material_ingest_permission(session, material, current_user)

    existing_result = await session.execute(
        select(KnowledgeDocument).where(
            KnowledgeDocument.source_type == KnowledgeSourceType.CLASS_MATERIAL,
            KnowledgeDocument.source_id == material.id,
        )
    )
    document = existing_result.scalar_one_or_none()
    if not document:
        document = KnowledgeDocument(
            source_type=KnowledgeSourceType.CLASS_MATERIAL,
            source_id=material.id,
            class_id=material.class_id,
            teacher_id=classroom.teacher_id,
            title=material.title,
            file_path=str(_resolve_material_file_path(material.file_url)),
            mime_type=material.file_type,
            parse_status=KnowledgeParseStatus.PENDING,
            metadata_json={"material_id": material.id},
        )
        session.add(document)
        await session.commit()
        await session.refresh(document)

    document.title = material.title
    document.file_path = str(_resolve_material_file_path(material.file_url))
    document.mime_type = material.file_type
    document.parse_status = KnowledgeParseStatus.PROCESSING
    document.parse_error = None
    document.updated_at = datetime.utcnow()
    session.add(document)
    await session.commit()

    try:
        file_path = _resolve_material_file_path(material.file_url)
        text, metadata = _extract_text_from_material(material, file_path)
        chunks = split_text(text)
        if not chunks:
            raise ValueError("No extractable text found")

        await session.execute(delete(KnowledgeChunk).where(KnowledgeChunk.document_id == document.id))
        for index, chunk in enumerate(chunks):
            session.add(
                KnowledgeChunk(
                    document_id=document.id,
                    chunk_index=index,
                    content=chunk,
                    content_type="theory",
                    token_count=estimate_token_count(chunk),
                    knowledge_tags=[],
                    metadata_json={
                        **metadata,
                        "class_id": material.class_id,
                        "document_name": material.title,
                    },
                    embedding=[],
                )
            )

        document.parse_status = KnowledgeParseStatus.COMPLETED
        document.parse_error = None
        document.metadata_json = {
            **document.metadata_json,
            **metadata,
            "chunk_count": len(chunks),
        }
    except Exception as exc:
        document.parse_status = KnowledgeParseStatus.FAILED
        document.parse_error = str(exc)

    document.updated_at = datetime.utcnow()
    session.add(document)
    await session.commit()
    await session.refresh(document)
    return document


@router.get("/search", response_model=KnowledgeSearchResponse)
async def search_knowledge_chunks(
    *,
    session: AsyncSession = Depends(get_session),
    current_user: User = Depends(deps.get_current_active_user),
    q: str,
    class_id: Optional[int] = None,
    source_type: Optional[KnowledgeSourceType] = None,
    skip: int = 0,
    limit: int = 20,
) -> Any:
    limit = max(1, min(limit, 100))
    filters = []
    if class_id is not None:
        filters.append(KnowledgeDocument.class_id == class_id)
    if source_type is not None:
        filters.append(KnowledgeDocument.source_type == source_type)
    if current_user.role != UserRole.ADMIN:
        filters.append(
            (KnowledgeDocument.class_id == None) |  # noqa: E711
            (KnowledgeDocument.teacher_id == current_user.id)
        )

    stmt = (
        select(KnowledgeChunk)
        .join(KnowledgeDocument, KnowledgeChunk.document_id == KnowledgeDocument.id)
        .where(KnowledgeChunk.content.ilike(f"%{q}%"), *filters)
        .order_by(KnowledgeChunk.created_at.desc())
    )
    count_stmt = select(func.count()).select_from(stmt.subquery())
    total = (await session.execute(count_stmt)).scalar_one()
    result = await session.execute(stmt.offset(skip).limit(limit))
    items = result.scalars().all()
    return KnowledgeSearchResponse(items=items, total=total)