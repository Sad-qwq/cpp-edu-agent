from datetime import datetime
from pathlib import Path
from typing import Any
from types import SimpleNamespace

from sqlalchemy.ext.asyncio import AsyncSession
from sqlmodel import delete, select

from app.ai.retrieval.chunking import estimate_token_count, split_text
from app.models.ai_question_generation import (
    KnowledgeChunk,
    KnowledgeDocument,
    KnowledgeParseStatus,
    KnowledgeSourceType,
)
from app.models.material import Material

BACKEND_ROOT = Path(__file__).resolve().parents[3]


class MaterialIngestionError(ValueError):
    pass


def resolve_material_file_path(file_url: str) -> Path:
    relative_path = file_url.replace("/static/", "uploads/").lstrip("/")
    return BACKEND_ROOT / relative_path


def extract_text_from_material(material: Material, file_path: Path) -> tuple[str, dict[str, Any]]:
    suffix = file_path.suffix.lower()
    metadata = {
        "material_id": material.id,
        "material_type": material.file_type,
        "extract_mode": "fallback",
    }

    if suffix in {".txt", ".md", ".py", ".cpp", ".c", ".h", ".hpp", ".json"} and file_path.exists():
        text = file_path.read_text(encoding="utf-8", errors="ignore")
        metadata["extract_mode"] = "direct_text"
        return text, metadata

    if suffix == ".pdf" and file_path.exists():
        try:
            from pypdf import PdfReader
        except ModuleNotFoundError as exc:
            raise MaterialIngestionError("PDF 自动入库依赖缺失，请安装 pypdf 后重试") from exc

        reader = PdfReader(str(file_path))
        text = "\n".join((page.extract_text() or "") for page in reader.pages)
        metadata["extract_mode"] = "pdf"
        return text, metadata

    if suffix == ".ppt" and file_path.exists():
        raise MaterialIngestionError("旧版 .ppt 文件暂不支持自动入库，请先转换为 .pptx 或 PDF 后再上传")

    if suffix == ".pptx" and file_path.exists():
        try:
            from pptx import Presentation
        except ModuleNotFoundError as exc:
            raise MaterialIngestionError("PPTX 自动入库依赖缺失，请安装 python-pptx 后重试") from exc

        try:
            presentation = Presentation(str(file_path))
        except Exception as exc:
            raise MaterialIngestionError(f"PPTX 解析失败，请确认文件未损坏并重试：{exc}") from exc

        texts: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        metadata["extract_mode"] = "ppt"
        return "\n".join(texts), metadata

    fallback_parts = [material.title]
    if material.description:
        fallback_parts.append(material.description)
    fallback_parts.append(f"资料类型: {material.file_type}")
    if material.file_url:
        fallback_parts.append(f"文件地址: {material.file_url}")
    return "\n".join(fallback_parts), metadata


async def upsert_knowledge_document(
    session: AsyncSession,
    *,
    source_type: KnowledgeSourceType,
    source_id: int | None,
    class_id: int | None,
    owner_id: int | None,
    title: str,
    description: str | None,
    file_url: str,
    file_type: str,
    extra_metadata: dict[str, Any] | None = None,
) -> KnowledgeDocument:
    existing_result = await session.execute(
        select(KnowledgeDocument).where(
            KnowledgeDocument.source_type == source_type,
            KnowledgeDocument.source_id == source_id,
        )
    )
    document = existing_result.scalar_one_or_none()
    if not document:
        document = KnowledgeDocument(
            source_type=source_type,
            source_id=source_id,
            class_id=class_id,
            teacher_id=owner_id,
            title=title,
            file_path=str(resolve_material_file_path(file_url)),
            mime_type=file_type,
            parse_status=KnowledgeParseStatus.PENDING,
            metadata_json=extra_metadata or {},
        )
        session.add(document)
        await session.flush()

    document.title = title
    document.file_path = str(resolve_material_file_path(file_url))
    document.mime_type = file_type
    document.class_id = class_id
    document.teacher_id = owner_id
    document.parse_status = KnowledgeParseStatus.PROCESSING
    document.parse_error = None
    document.updated_at = datetime.utcnow()
    if extra_metadata is not None:
        document.metadata_json = {**document.metadata_json, **extra_metadata}
    session.add(document)
    await session.flush()

    source = SimpleNamespace(
        id=source_id,
        title=title,
        description=description,
        file_url=file_url,
        file_type=file_type,
        class_id=class_id,
    )

    try:
        file_path = resolve_material_file_path(file_url)
        text, metadata = extract_text_from_material(source, file_path)
        chunks = split_text(text)
        if not chunks:
            raise ValueError("No extractable text found")

        await session.execute(delete(KnowledgeChunk).where(KnowledgeChunk.document_id == document.id))
        for index, chunk in enumerate(chunks):
            session.add(
                KnowledgeChunk(
                    document_id=document.id,
                    chunk_index=index,
                    content=chunk,
                    content_type="theory",
                    token_count=estimate_token_count(chunk),
                    knowledge_tags=[],
                    metadata_json={
                        **metadata,
                        **(extra_metadata or {}),
                        "class_id": class_id,
                        "document_name": title,
                    },
                    embedding=[],
                )
            )

        document.parse_status = KnowledgeParseStatus.COMPLETED
        document.parse_error = None
        document.metadata_json = {
            **document.metadata_json,
            **metadata,
            **(extra_metadata or {}),
            "chunk_count": len(chunks),
        }
    except Exception as exc:
        document.parse_status = KnowledgeParseStatus.FAILED
        document.parse_error = str(exc)

    document.updated_at = datetime.utcnow()
    session.add(document)
    await session.flush()
    return document


async def upsert_material_knowledge_document(
    session: AsyncSession,
    *,
    material: Material,
    teacher_id: int,
) -> KnowledgeDocument:
    return await upsert_knowledge_document(
        session,
        source_type=KnowledgeSourceType.CLASS_MATERIAL,
        source_id=material.id,
        class_id=material.class_id,
        owner_id=teacher_id,
        title=material.title,
        description=material.description,
        file_url=material.file_url,
        file_type=material.file_type,
        extra_metadata={"material_id": material.id},
    )


async def delete_material_knowledge_documents(session: AsyncSession, *, material_id: int) -> None:
    result = await session.execute(
        select(KnowledgeDocument.id).where(
            KnowledgeDocument.source_type == KnowledgeSourceType.CLASS_MATERIAL,
            KnowledgeDocument.source_id == material_id,
        )
    )
    document_ids = result.scalars().all()
    if not document_ids:
        return

    await session.execute(delete(KnowledgeChunk).where(KnowledgeChunk.document_id.in_(document_ids)))
    await session.execute(delete(KnowledgeDocument).where(KnowledgeDocument.id.in_(document_ids)))


async def delete_knowledge_document(session: AsyncSession, *, document_id: int) -> None:
    await session.execute(delete(KnowledgeChunk).where(KnowledgeChunk.document_id == document_id))
    await session.execute(delete(KnowledgeDocument).where(KnowledgeDocument.id == document_id))